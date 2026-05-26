"""
Genera el executive summary del hackathon: 5 slides 16:9.

Salidas en submission/:
  - executive_summary.pptx  (vidos Vision/PPCF/DOS embebidos en slides 3-5)
  - executive_summary.pdf   (version estatica, los videos salen como still)

Slides: 1 intro diagonalidad, 2 pipeline, 3 Vision, 4 PPCF, 5 DOS.
Ejecutar desde la raiz del repo:  python3 deliverable/make_exec_summary.py
"""
from pathlib import Path
import sys
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))
import src.viz.common as _vc        # registra Chakra Petch + rcParams base   # noqa: E402

FIG  = ROOT / "deliverable" / "figures"
VID  = ROOT / "figures" / "videos"
LOGO = ROOT / "figures" / "logos" / "jo_logo.png"     # CCV-style JO logo
OUT  = ROOT / "submission"
SLD  = FIG / "exec_slides"
OUT.mkdir(exist_ok=True)
SLD.mkdir(exist_ok=True)

# --- Paleta LIGHT OPTA (CCV identity) ---------------------------------
BG      = "#ffffff"   # blanco puro
WHITE   = "#000000"   # alias retro-compat: codigo abajo usa WHITE para texto principal -> negro
MUTE    = "#666666"   # leyenda / muted text
ACCENT  = "#3b82f6"   # corporate blue (mismo ATT del common.py)

# Chakra Petch ya esta registrado al importar src.viz.common; aqui solo
# fijamos la familia activa.
plt.rcParams["font.family"] = "sans-serif"
plt.rcParams["font.sans-serif"] = _vc.FONT_STACK

# slide en unidades de grid 16 x 9
GX, GY = 16.0, 9.0
SLIDE_W_IN, SLIDE_H_IN = 13.333, 7.5


def new_slide():
    fig = plt.figure(figsize=(SLIDE_W_IN, SLIDE_H_IN), dpi=200)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, GX); ax.set_ylim(0, GY)
    ax.axis("off")
    ax.add_patch(plt.Rectangle((0, 0), GX, GY, color=BG, zorder=0))
    # banda de acento superior
    ax.add_patch(plt.Rectangle((0, GY - 0.09), GX, 0.09, color=ACCENT, zorder=1))
    return fig, ax


def footer(ax, page):
    ax.text(0.55, 0.42, "Diagonality · The Best of Both Worlds",
            color=MUTE, fontsize=10, va="center", zorder=6)
    ax.text(GX - 0.55, 0.42, f"{page} / 5", color=MUTE, fontsize=10,
            ha="right", va="center", zorder=6)
    if LOGO.exists():
        img = plt.imread(str(LOGO))
        ar = img.shape[1] / img.shape[0]
        h = 0.62; w = h * ar
        ax.imshow(img, extent=[GX - 0.45 - w, GX - 0.45, GY - 0.88, GY - 0.88 + h],
                  zorder=6, aspect="auto")


def place_image(ax, path, x0, x1, y0, y1, z=5):
    """Encaja la imagen en la caja preservando aspecto. Devuelve el rect real."""
    img = plt.imread(str(path))
    ar_img = img.shape[1] / img.shape[0]
    bw, bh = x1 - x0, y1 - y0
    if ar_img > bw / bh:
        w = bw; h = w / ar_img
    else:
        h = bh; w = h * ar_img
    cx, cy = (x0 + x1) / 2, (y0 + y1) / 2
    ex = [cx - w / 2, cx + w / 2, cy - h / 2, cy + h / 2]
    ax.imshow(img, extent=ex, zorder=z, aspect="auto")
    ax.add_patch(plt.Rectangle((ex[0], ex[2]), ex[1] - ex[0], ex[3] - ex[2],
                 fill=False, ec="#dddddd", lw=1.0, zorder=z + 1))
    return ex  # [left, right, bottom, top]


def kicker(ax, text):
    ax.text(0.7, GY - 0.62, text.upper(), color=ACCENT, fontsize=12.5,
            fontweight="bold", va="center", zorder=6)


def title(ax, text, y=GY - 1.35):
    ax.text(0.7, y, text, color=WHITE, fontsize=29, fontweight="bold",
            va="center", zorder=6)


def save(fig, name):
    p = SLD / name
    fig.savefig(str(p), facecolor=BG)
    plt.close(fig)
    return p


# ====================================================================
# SLIDE 1 -- portada minimalista: titulo + una linea de que es
# ====================================================================
def slide1():
    fig, ax = new_slide()
    cx = GX / 2

    # Kicker centrado
    ax.text(cx, 7.55,
            "AWS WORLD SPORTS INNOVATION CUP 2026     ·     CHALLENGE 2",
            color=ACCENT, fontsize=12.5, fontweight="bold",
            ha="center", va="center", zorder=6)

    # Titulo en dos niveles
    ax.text(cx, 5.85, "Diagonality", color=WHITE, fontsize=56,
            fontweight="bold", ha="center", va="center", zorder=6)
    ax.text(cx, 4.74, "The Best of Both Worlds", color=ACCENT,
            fontsize=26, ha="center", va="center", zorder=6)

    # Regla fina de acento
    ax.add_patch(plt.Rectangle((cx - 1.5, 4.10), 3.0, 0.03,
                               color=ACCENT, zorder=6))

    # La teoria: el blog Spielverlagerung y su tesis central
    ax.text(cx, 3.42,
            "The tactical blog Spielverlagerung argues the diagonal is football's optimum:",
            color=WHITE, fontsize=13.5, ha="center", va="center", zorder=6)
    ax.text(cx, 2.98,
            "the progression of a vertical pass with the safety of a horizontal one.",
            color=WHITE, fontsize=13.5, ha="center", va="center", zorder=6)

    # Que hacemos nosotros
    ax.text(cx, 2.22,
            "We put that theory to the test on 3D skeleton data,",
            color=WHITE, fontsize=13.5, ha="center", va="center", zorder=6)
    ax.text(cx, 1.78,
            "and turn it into a live map of where to attack.",
            color=WHITE, fontsize=13.5, ha="center", va="center", zorder=6)

    # Autor
    ax.text(cx, 0.92, "Jaime Oriol", color=WHITE, fontsize=13,
            ha="center", va="center", zorder=6)

    footer(ax, 1)
    return save(fig, "slide1.png"), None


# ====================================================================
# SLIDE 2 -- la infografia del pipeline, a sangre (solo el PNG)
# ====================================================================
def slide2():
    fig, ax = new_slide()
    img = plt.imread(str(FIG / "pipeline.png"))
    # pipeline.png es 16:9 -> ocupa la slide entera, sin texto ni chrome.
    ax.imshow(img, extent=[0, GX, 0, GY], zorder=5, aspect="auto")
    return save(fig, "slide2.png"), None


# ====================================================================
# SLIDES 3-5 -- video stages
# ====================================================================
def video_slide(idx, kick, ttl, caption, fig_png, vid_mp4):
    fig, ax = new_slide()
    kicker(ax, kick)
    title(ax, ttl)
    rect = place_image(ax, FIG / fig_png, 0.7, 15.3, 1.5, GY - 2.05)
    ax.text(8.0, 1.05, caption, color=WHITE, fontsize=12.2, ha="center",
            va="center", zorder=6)
    footer(ax, idx)
    png = save(fig, f"slide{idx}.png")
    return png, {"rect": rect, "video": VID / vid_mp4, "poster": FIG / fig_png}


# ====================================================================
# build
# ====================================================================
def grid_to_pptx(rect):
    """rect = [left, right, bottom, top] en grid 16x9 -> (left,top,w,h) pulgadas."""
    sx = SLIDE_W_IN / GX
    sy = SLIDE_H_IN / GY
    left = rect[0] * sx
    width = (rect[1] - rect[0]) * sx
    top = (GY - rect[3]) * sy
    height = (rect[3] - rect[2]) * sy
    return Inches(left), Inches(top), Inches(width), Inches(height)


def main():
    slides = [
        slide1(),
        slide2(),
        video_slide(3, "Stage 1 · Vision",
                    "What every player can see",
                    "A 120° cone of vision built from where each player is really "
                    "looking. It narrows when they sprint, and other players block the view.",
                    "Vision.png", "Vision_Video.mp4"),
        video_slide(4, "Stage 2 · Pitch Control",
                    "What every player can reach",
                    "How much ground a player covers depends on which way he faces. "
                    "Behind his shoulder it shrinks, leaving a real gap in the defence.",
                    "PPCF.png", "PPCF_Video.mp4"),
        video_slide(5, "Stage 3 · DOS",
                    "The Diagonal Opportunity Surface",
                    "The extra space a diagonal pass opens up, shown only where the player "
                    "on the ball can actually see it. Cyan: he sees it. Amber: he lost track of it.",
                    "DOS.png", "DOS_Video.mp4"),
    ]

    # --- PDF (5 paginas, estatico) ------------------------------------
    imgs = [Image.open(str(p)).convert("RGB") for p, _ in slides]
    pdf = OUT / "executive_summary.pdf"
    imgs[0].save(str(pdf), save_all=True, append_images=imgs[1:])
    print("PDF  ->", pdf)

    # --- PPTX (videos embebidos en 3-5) -------------------------------
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]
    for png, vid in slides:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(str(png), 0, 0, width=prs.slide_width,
                             height=prs.slide_height)
        if vid is not None:
            left, top, w, h = grid_to_pptx(vid["rect"])
            s.shapes.add_movie(str(vid["video"]), left, top, w, h,
                               poster_frame_image=str(vid["poster"]),
                               mime_type="video/mp4")
    pptx = OUT / "executive_summary.pptx"
    prs.save(str(pptx))
    print("PPTX ->", pptx)


if __name__ == "__main__":
    main()
